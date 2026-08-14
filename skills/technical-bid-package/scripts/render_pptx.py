#!/usr/bin/env python3
import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is required: pip install python-pptx", file=sys.stderr)
    raise


def inspect_template(path: Path) -> int:
    presentation = Presentation(path)
    result = []
    for index, layout in enumerate(presentation.slide_layouts):
        placeholders = [{"index": item.placeholder_format.idx, "name": item.name, "type": str(item.placeholder_format.type)} for item in layout.placeholders]
        result.append({"index": index, "name": layout.name, "placeholders": placeholders})
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def set_placeholder(slide, index, text, bullets=False):
    placeholder = next((item for item in slide.placeholders if item.placeholder_format.idx == index), None)
    if placeholder is None or not hasattr(placeholder, "text_frame"):
        return False
    frame = placeholder.text_frame
    frame.clear()
    if bullets:
        for position, value in enumerate(text):
            paragraph = frame.paragraphs[0] if position == 0 else frame.add_paragraph()
            paragraph.text = str(value)
            paragraph.level = 0
    else:
        frame.text = str(text)
    return True


def render(template: Path, plan_path: Path, map_path: Path, output: Path) -> int:
    presentation = Presentation(template)
    plan = json.loads(plan_path.read_text(encoding="utf-8-sig"))
    layout_map = json.loads(map_path.read_text(encoding="utf-8-sig"))
    layouts = {layout.name: layout for layout in presentation.slide_layouts}
    warnings = []

    title_config = layout_map.get("title", {})
    title_layout = layouts.get(title_config.get("layout"), presentation.slide_layouts[0])
    slide = presentation.slides.add_slide(title_layout)
    if not set_placeholder(slide, title_config.get("titlePlaceholder", 0), plan.get("title", "")):
        warnings.append("title placeholder not found on title slide")
    set_placeholder(slide, title_config.get("bodyPlaceholder", 1), plan.get("subtitle", ""))

    for item in plan.get("slides", []):
        role = item.get("role", "content")
        config = layout_map.get(role, layout_map.get("content", {}))
        layout = layouts.get(config.get("layout"))
        if layout is None:
            layout = presentation.slide_layouts[min(1, len(presentation.slide_layouts) - 1)]
            warnings.append(f"layout not found for role {role}; fallback used")
        slide = presentation.slides.add_slide(layout)
        if not set_placeholder(slide, config.get("titlePlaceholder", 0), item.get("title", "")):
            warnings.append(f"title placeholder not found for slide {item.get('title', '<untitled>')}")
        if not set_placeholder(slide, config.get("bodyPlaceholder", 1), item.get("bullets", []), bullets=True):
            warnings.append(f"body placeholder not found for slide {item.get('title', '<untitled>')}")

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    print(json.dumps({"ok": True, "output": str(output), "slides": len(presentation.slides), "warnings": warnings}, ensure_ascii=False, indent=2))
    return 0


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--inspect-template", type=Path)
    parser.add_argument("--template", type=Path)
    parser.add_argument("--plan", type=Path)
    parser.add_argument("--map", dest="map_path", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()
    if args.inspect_template:
        return inspect_template(args.inspect_template)
    if not all([args.template, args.plan, args.map_path, args.output]):
        parser.error("--template, --plan, --map, and --output are required for rendering")
    return render(args.template, args.plan, args.map_path, args.output)


if __name__ == "__main__":
    sys.exit(main())
